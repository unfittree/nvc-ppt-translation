from __future__ import annotations

import csv
import re
from dataclasses import dataclass
from pathlib import Path
import subprocess
import tempfile

from deep_translator import GoogleTranslator
from hanziconv import HanziConv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


_HAS_LATIN_TEXT = re.compile(r"[A-Za-z]")
_SOURCE_WORD_CHARS = re.compile(r"[A-Za-z]")
_UNRESOLVED_GLOSSARY_TOKEN = re.compile(r"\[\[\[\d+]]]")


@dataclass(frozen=True)
class GlossaryEntry:
    pattern: re.Pattern[str]
    target: str
    priority: int


_BUILTIN_GLOSSARY: tuple[tuple[str, str], ...] = (
    (r"\bnon[\s-]?violent communication\b", "非暴力沟通"),
    (r"\bcompassionate communication\b", "同理沟通"),
    (r"\bNVC\b", "NVC（非暴力沟通）"),
    (r"\bself[-\s]?empathy\b", "自我同理"),
    (r"\bempathy\b", "同理心"),
    (r"\bobservation(s)?\b", "观察"),
    (r"\bfeeling(s)?\b", "感受"),
    (r"\bneed(s)?\b", "需要"),
    (r"\brequest(s)?\b", "请求"),
    (r"\bdemand(s|ed|ing)?\b", "要求"),
    (r"\bjudgment(s|al)?\b", "评判"),
    (r"\bblame(s|d)?\b", "指责"),
    (r"\bcriticism\b", "批评"),
    (r"\bself[-\s]?awareness\b", "自我觉察"),
    (r"\binner growth\b", "内在成长"),
    (r"\binner child\b", "内在小孩"),
    (r"\btrigger(s|ed|ing)?\b", "触发点"),
    (r"\bboundar(y|ies)\b", "边界"),
    (r"\bvulnerability\b", "脆弱"),
    (r"\bauthenticity\b", "真实"),
    (r"\bmindfulness\b", "正念"),
    (r"\bcompassion\b", "慈悲"),
    (r"\bself[-\s]?compassion\b", "自我慈悲"),
    (r"\bconnection\b", "连接"),
    (r"\bdisconnection\b", "失联"),
    (r"\bshame\b", "羞耻感"),
    (r"\bguilt\b", "内疚"),
    (r"\banger\b", "愤怒"),
    (r"\bgratitude\b", "感恩"),
    (r"\bconflict\b", "冲突"),
    (r"\breconciliation\b", "和解"),
    (r"\bintimacy\b", "亲密"),
    (r"\bcouple(s)?\b", "伴侣"),
    (r"\bpartnership\b", "亲密关系"),
    (r"\bgiraffe language\b", "长颈鹿语言"),
    (r"\bjackal language\b", "豺狗语言"),
)


@dataclass
class TranslationStats:
    slides: int = 0
    text_segments_seen: int = 0
    text_segments_translated: int = 0
    cache_hits: int = 0
    glossary_replacements: int = 0


def _build_builtin_glossary() -> list[GlossaryEntry]:
    entries: list[GlossaryEntry] = []
    for pattern, target in _BUILTIN_GLOSSARY:
        entries.append(
            GlossaryEntry(
                pattern=re.compile(pattern, re.IGNORECASE),
                target=target,
                priority=len(pattern),
            )
        )
    return entries


def _build_literal_glossary_pattern(source_term: str) -> re.Pattern[str]:
    escaped = re.escape(source_term)
    if _SOURCE_WORD_CHARS.search(source_term):
        return re.compile(rf"(?<![A-Za-z]){escaped}(?![A-Za-z])", re.IGNORECASE)
    return re.compile(escaped, re.IGNORECASE)


def load_glossary_file(glossary_file: str | Path | None) -> list[GlossaryEntry]:
    if not glossary_file:
        return []

    path = Path(glossary_file)
    if not path.exists():
        raise FileNotFoundError(f"Glossary file not found: {path}")

    entries: list[GlossaryEntry] = []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            if not row:
                continue

            first_cell = row[0].strip()
            if not first_cell or first_cell.startswith("#"):
                continue

            if len(row) < 2:
                continue

            source = first_cell
            target = row[1].strip()
            if not target:
                continue

            entries.append(
                GlossaryEntry(
                    pattern=_build_literal_glossary_pattern(source),
                    target=target,
                    priority=len(source),
                )
            )
    return entries


class TerminologyGlossary:
    def __init__(self, extra_entries: list[GlossaryEntry] | None = None) -> None:
        entries = _build_builtin_glossary()
        if extra_entries:
            entries.extend(extra_entries)
        self._entries = sorted(entries, key=lambda entry: entry.priority, reverse=True)

    def protect(self, text: str) -> tuple[str, dict[str, str], int]:
        protected = text
        token_index = 0
        token_map: dict[str, str] = {}
        hits = 0

        for entry in self._entries:
            def _replacement(_: re.Match[str]) -> str:
                nonlocal token_index, hits
                token = f"[[[{token_index}]]]"
                token_index += 1
                token_map[token] = entry.target
                hits += 1
                return token

            protected = entry.pattern.sub(_replacement, protected)

        return protected, token_map, hits

    @staticmethod
    def restore(text: str, token_map: dict[str, str]) -> str:
        restored = text
        for token, target in token_map.items():
            restored = restored.replace(token, target)

        # Defensive fallback in case spacing is introduced inside token brackets.
        index_to_target: dict[str, str] = {}
        for token, target in token_map.items():
            match = re.fullmatch(r"\[\[\[(\d+)]]]", token)
            if match:
                index_to_target[match.group(1)] = target
        if index_to_target:
            token_like = re.compile(r"\[\s*\[\s*\[\s*(\d+)\s*\]\s*\]\s*\]")
            restored = token_like.sub(
                lambda match: index_to_target.get(match.group(1), match.group(0)),
                restored,
            )
        return restored


def _split_by_weights(text: str, weights: list[int]) -> list[str]:
    if not weights:
        return []

    total_weight = sum(weights)
    if total_weight <= 0:
        return [""] * len(weights)

    total_chars = len(text)
    raw_allocations = [(weight * total_chars) / total_weight for weight in weights]
    int_allocations = [int(allocation) for allocation in raw_allocations]
    assigned = sum(int_allocations)
    remainder = total_chars - assigned

    if remainder > 0:
        order = sorted(
            range(len(weights)),
            key=lambda index: raw_allocations[index] - int_allocations[index],
            reverse=True,
        )
        for index in order[:remainder]:
            int_allocations[index] += 1

    chunks: list[str] = []
    cursor = 0
    for allocation in int_allocations:
        next_cursor = cursor + allocation
        chunks.append(text[cursor:next_cursor])
        cursor = next_cursor

    if cursor < total_chars:
        chunks[-1] += text[cursor:]

    return chunks


def _convert_ppt_to_pptx(input_path: Path, converted_path: Path) -> None:
    escaped_input = str(input_path).replace("'", "''")
    escaped_output = str(converted_path).replace("'", "''")
    script = (
        "$ErrorActionPreference='Stop';"
        f"$inputPath='{escaped_input}';"
        f"$outputPath='{escaped_output}';"
        "$ppt=$null;$pres=$null;"
        "try {"
        "$ppt=New-Object -ComObject PowerPoint.Application;"
        "$ppt.Visible=1;"
        "$pres=$ppt.Presentations.Open($inputPath,$false,$false,$false);"
        "$pres.SaveAs($outputPath,24);"
        "} finally {"
        "if ($pres -ne $null) { $pres.Close() };"
        "if ($ppt -ne $null) { $ppt.Quit() }"
        "}"
    )
    result = subprocess.run(
        ["powershell", "-NoProfile", "-Command", script],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0 or not converted_path.exists():
        stderr = result.stderr.strip() or "Unknown PowerPoint conversion error."
        raise RuntimeError(
            "Failed to convert .ppt to .pptx via PowerPoint automation. "
            f"Details: {stderr}"
        )


class CachedGoogleTranslator:
    def __init__(
        self,
        source_lang: str = "en",
        target_lang: str = "zh-CN",
        glossary: TerminologyGlossary | None = None,
        force_simplified: bool = True,
    ) -> None:
        self._translator = GoogleTranslator(source=source_lang, target=target_lang)
        self._cache: dict[str, str] = {}
        self._glossary = glossary or TerminologyGlossary()
        self._force_simplified = force_simplified and target_lang.lower().startswith("zh")

    def _normalize_target_text(self, text: str) -> str:
        if self._force_simplified:
            return HanziConv.toSimplified(text)
        return text

    def translate(self, text: str) -> tuple[str, bool, int]:
        stripped = text.strip()
        if not stripped:
            return text, False, 0

        if not _HAS_LATIN_TEXT.search(text):
            return self._normalize_target_text(text), False, 0

        cached = self._cache.get(text)
        if cached is not None:
            return cached, True, 0

        protected, token_map, glossary_hits = self._glossary.protect(text)
        translated = self._translator.translate(protected)
        translated = self._glossary.restore(translated, token_map)
        translated = self._normalize_target_text(translated)
        if _UNRESOLVED_GLOSSARY_TOKEN.search(translated):
            raise RuntimeError(
                "Unresolved glossary token detected after translation; aborted to avoid placeholder leakage."
            )
        self._cache[text] = translated
        return translated, False, glossary_hits


def _translate_text_frame(text_frame, translator: CachedGoogleTranslator, stats: TranslationStats) -> None:
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            original_runs = [run.text or "" for run in paragraph.runs]
            original_paragraph_text = "".join(original_runs)
            if not original_paragraph_text:
                continue

            stats.text_segments_seen += 1
            translated, from_cache, glossary_hits = translator.translate(original_paragraph_text)
            if from_cache:
                stats.cache_hits += 1
            stats.glossary_replacements += glossary_hits

            if translated != original_paragraph_text:
                weights = [len(text) for text in original_runs]
                translated_runs = _split_by_weights(translated, weights)
                for run, run_text in zip(paragraph.runs, translated_runs):
                    run.text = run_text
                stats.text_segments_translated += 1
        else:
            original = paragraph.text
            if not original:
                continue

            stats.text_segments_seen += 1
            translated, from_cache, glossary_hits = translator.translate(original)
            if from_cache:
                stats.cache_hits += 1
            stats.glossary_replacements += glossary_hits

            if translated != original:
                paragraph.text = translated
                stats.text_segments_translated += 1


def _translate_shape(shape, translator: CachedGoogleTranslator, stats: TranslationStats) -> None:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for grouped_shape in shape.shapes:
            _translate_shape(grouped_shape, translator, stats)
        return

    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        _translate_text_frame(shape.text_frame, translator, stats)

    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _translate_text_frame(cell.text_frame, translator, stats)


def translate_presentation(
    input_path: str | Path,
    output_path: str | Path,
    source_lang: str = "en",
    target_lang: str = "zh-CN",
    include_notes: bool = False,
    glossary_file: str | Path | None = None,
    force_simplified: bool = True,
) -> TranslationStats:
    source_path = Path(input_path)
    if source_path.suffix.lower() not in {".pptx", ".ppt"}:
        raise ValueError("Input file must be .pptx or .ppt")

    if Path(output_path).suffix.lower() != ".pptx":
        raise ValueError("Output file must be .pptx")

    extra_entries = load_glossary_file(glossary_file)
    translator = CachedGoogleTranslator(
        source_lang=source_lang,
        target_lang=target_lang,
        glossary=TerminologyGlossary(extra_entries=extra_entries),
        force_simplified=force_simplified,
    )

    with tempfile.TemporaryDirectory(prefix="ppt_translation_") as temp_dir:
        presentation_input = source_path
        if source_path.suffix.lower() == ".ppt":
            converted = Path(temp_dir) / f"{source_path.stem}.pptx"
            _convert_ppt_to_pptx(source_path, converted)
            presentation_input = converted

        presentation = Presentation(str(presentation_input))
        stats = TranslationStats(slides=len(presentation.slides))

        for slide in presentation.slides:
            for shape in slide.shapes:
                _translate_shape(shape, translator, stats)

            if include_notes and slide.has_notes_slide:
                _translate_text_frame(slide.notes_slide.notes_text_frame, translator, stats)

        presentation.save(str(output_path))

    return stats
